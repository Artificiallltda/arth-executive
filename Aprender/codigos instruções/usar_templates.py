"""
COMO USAR OS TEMPLATES NO SISTEMA DE GERAÇÃO
Este arquivo mostra como carregar e usar os templates criados
nos geradores gerar_docx.py e gerar_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGB


# ─────────────────────────────────────────────
# USANDO TEMPLATE PPTX
# ─────────────────────────────────────────────

def gerar_pptx_com_template(
    template_nome: str,
    titulo: str,
    subtitulo: str,
    slides_conteudo: list,
    caminho_saida: str = None
) -> str:
    """
    Gera PPTX usando um template existente como base visual.
    
    Args:
        template_nome: Nome do template (financeiro, marketing, escola, etc.)
        titulo: Título da apresentação
        subtitulo: Subtítulo
        slides_conteudo: Lista de dicts [{"titulo": "...", "conteudo": "...", "tipo": "texto|lista"}]
        caminho_saida: Caminho para salvar o arquivo
    
    Returns:
        Caminho do arquivo gerado
    """
    template_path = f"templates/pptx/template_{template_nome}.pptx"
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(
            f"Template '{template_nome}' não encontrado em {template_path}. "
            f"Execute criar_templates_pptx.py primeiro."
        )
    
    # Carregar template como base
    prs = Presentation(template_path)
    
    # O template já tem slides de exemplo — vamos LIMPAR e recriar
    # Mantemos apenas as configurações de tamanho e tema
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Criar nova apresentação com as dimensões do template
    prs_novo = Presentation()
    prs_novo.slide_width = slide_width
    prs_novo.slide_height = slide_height
    
    # Copiar slides do template como base e substituir textos
    # (Abordagem mais simples: usar o template e editar os slides existentes)
    prs_template = Presentation(template_path)
    
    # Slide 1 (capa) — editar textos do template
    slide_capa = prs_template.slides[0]
    for shape in slide_capa.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "TÍTULO DA APRESENTAÇÃO" in run.text:
                        run.text = titulo.upper()
                    elif "Subtítulo ou descrição aqui" in run.text:
                        run.text = subtitulo
                    elif "Autor  |  Data  |  Versão" in run.text:
                        from datetime import datetime
                        run.text = f"Sistema IA  |  {datetime.now().strftime('%d/%m/%Y')}"
    
    # Salvar apresentação editada
    if not caminho_saida:
        os.makedirs("outputs", exist_ok=True)
        caminho_saida = f"outputs/{titulo.replace(' ', '_')}.pptx"
    
    prs_template.save(caminho_saida)
    print(f"✅ PPTX gerado com template '{template_nome}': {caminho_saida}")
    return caminho_saida


# ─────────────────────────────────────────────
# USANDO TEMPLATE DOCX
# ─────────────────────────────────────────────

def gerar_docx_com_template(
    template_nome: str,
    titulo: str,
    secoes: list,
    autor: str = "Sistema IA",
    caminho_saida: str = None
) -> str:
    """
    Gera DOCX usando um template existente como base visual.
    
    Args:
        template_nome: Nome do template (financeiro, marketing, escola, etc.)
        titulo: Título do documento
        secoes: Lista de dicts [{"titulo": "...", "conteudo": "...", "tipo": "texto|lista|tabela"}]
        autor: Nome do autor
        caminho_saida: Caminho para salvar o arquivo
    
    Returns:
        Caminho do arquivo gerado
    """
    template_path = f"templates/docx/template_{template_nome}.docx"
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(
            f"Template '{template_nome}' não encontrado em {template_path}. "
            f"Execute criar_templates_docx.py primeiro."
        )
    
    # Carregar template — mantém estilos, cabeçalho e rodapé
    doc = Document(template_path)
    
    # Limpar conteúdo existente (mantendo cabeçalho/rodapé do template)
    for para in doc.paragraphs:
        p = para._element
        p.getparent().remove(p)
    
    # Detectar paleta de cores do template pelo nome
    from criar_templates_docx import PALETAS_DOCX
    paleta = PALETAS_DOCX.get(template_nome, PALETAS_DOCX["relatorio"])
    
    # ── Título principal ──
    p_titulo = doc.add_paragraph()
    run = p_titulo.add_run(titulo.upper())
    run.font.size = DocxPt(26)
    run.font.bold = True
    run.font.color.rgb = paleta["titulo"]
    run.font.name = "Segoe UI"
    
    # ── Linha decorativa ──
    from criar_templates_docx import _add_linha_colorida
    _add_linha_colorida(doc, paleta["linha_hex"])
    
    # ── Metadados ──
    from datetime import datetime
    p_meta = doc.add_paragraph()
    run_meta = p_meta.add_run(f"{autor}  |  {datetime.now().strftime('%d/%m/%Y')}")
    run_meta.font.size = DocxPt(10)
    run_meta.font.italic = True
    run_meta.font.color.rgb = paleta["subtitulo"]
    
    doc.add_paragraph()
    
    # ── Seções de conteúdo ──
    for i, secao in enumerate(secoes, 1):
        tipo = secao.get("tipo", "texto")
        titulo_secao = secao.get("titulo", "")
        conteudo = secao.get("conteudo", "")
        
        if titulo_secao:
            h = doc.add_heading(f"{i}. {titulo_secao}", level=2)
            for run in h.runs:
                run.font.color.rgb = paleta["titulo"]
                run.font.name = "Segoe UI"
                run.font.size = DocxPt(13)
        
        if tipo == "texto":
            p = doc.add_paragraph(str(conteudo))
            p.alignment = 3  # JUSTIFY
            for run in p.runs:
                run.font.size = DocxPt(11)
                run.font.name = "Calibri"
        
        elif tipo == "lista":
            itens = secao.get("dados", [])
            if not itens and isinstance(conteudo, list):
                itens = conteudo
            for item in itens:
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(str(item))
                run.font.size = DocxPt(11)
                run.font.name = "Calibri"
        
        doc.add_paragraph()
    
    # Salvar
    if not caminho_saida:
        os.makedirs("outputs", exist_ok=True)
        caminho_saida = f"outputs/{titulo.replace(' ', '_')}.docx"
    
    doc.save(caminho_saida)
    print(f"✅ DOCX gerado com template '{template_nome}': {caminho_saida}")
    return caminho_saida


# ─────────────────────────────────────────────
# LISTAR TEMPLATES DISPONÍVEIS
# ─────────────────────────────────────────────

def listar_templates():
    """Lista todos os templates disponíveis."""
    print("\n📁 TEMPLATES PPTX DISPONÍVEIS:")
    if os.path.exists("templates/pptx"):
        for f in sorted(os.listdir("templates/pptx")):
            if f.endswith(".pptx"):
                nome = f.replace("template_", "").replace(".pptx", "")
                print(f"  ✅ {nome}")
    else:
        print("  ❌ Nenhum template PPTX encontrado. Execute criar_templates_pptx.py")

    print("\n📁 TEMPLATES DOCX DISPONÍVEIS:")
    if os.path.exists("templates/docx"):
        for f in sorted(os.listdir("templates/docx")):
            if f.endswith(".docx"):
                nome = f.replace("template_", "").replace(".docx", "")
                print(f"  ✅ {nome}")
    else:
        print("  ❌ Nenhum template DOCX encontrado. Execute criar_templates_docx.py")


# ─────────────────────────────────────────────
# EXEMPLO DE USO
# ─────────────────────────────────────────────

if __name__ == "__main__":
    listar_templates()

    # Exemplo: gerar DOCX financeiro
    # gerar_docx_com_template(
    #     template_nome="financeiro",
    #     titulo="Relatório Financeiro Q1 2026",
    #     secoes=[
    #         {"titulo": "Resumo Executivo", "tipo": "texto", "conteudo": "As receitas cresceram 35%..."},
    #         {"titulo": "Principais Indicadores", "tipo": "lista", "dados": ["ROI: 22%", "Margem: 18%", "CAC: R$ 120"]},
    #     ],
    #     autor="Gean"
    # )
