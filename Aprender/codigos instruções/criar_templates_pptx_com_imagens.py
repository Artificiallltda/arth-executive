"""
GERADOR DE PPTX COM IMAGENS CONTEXTUAIS
Fluxo completo: Gemini cria prompt → DALL-E gera imagem → imagem entra no slide
"""

import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from gerar_imagem_contextual import gerar_imagem_contextual

os.makedirs("outputs", exist_ok=True)

# ── Paletas (mesmas do template base) ──
PALETAS = {
    "financeiro":   {"fundo": RGBColor(0x0A,0x1F,0x44), "primaria": RGBColor(0x00,0x8B,0x8B), "secundaria": RGBColor(0x00,0xD4,0xAA), "acento": RGBColor(0xF0,0xB4,0x29), "texto": RGBColor(0xFF,0xFF,0xFF), "texto_sub": RGBColor(0xB0,0xC4,0xDE)},
    "marketing":    {"fundo": RGBColor(0x1A,0x00,0x2E), "primaria": RGBColor(0xFF,0x00,0x80), "secundaria": RGBColor(0xFF,0x6B,0x35), "acento": RGBColor(0xFF,0xD7,0x00), "texto": RGBColor(0xFF,0xFF,0xFF), "texto_sub": RGBColor(0xE0,0xC0,0xFF)},
    "apresentacao": {"fundo": RGBColor(0xFF,0xFF,0xFF), "primaria": RGBColor(0x2C,0x3E,0x50), "secundaria": RGBColor(0x3B,0x82,0xF6), "acento": RGBColor(0xE8,0xF4,0xFD), "texto": RGBColor(0x1E,0x29,0x3B), "texto_sub": RGBColor(0x64,0x74,0x8B)},
    "escola":       {"fundo": RGBColor(0xF0,0xF9,0xFF), "primaria": RGBColor(0x05,0x96,0x68), "secundaria": RGBColor(0xF5,0x9E,0x0B), "acento": RGBColor(0xFE,0xF3,0xC7), "texto": RGBColor(0x1E,0x29,0x3B), "texto_sub": RGBColor(0x37,0x41,0x51)},
    "ideias":       {"fundo": RGBColor(0xFF,0xFD,0xF0), "primaria": RGBColor(0xD9,0x77,0x06), "secundaria": RGBColor(0x7C,0x3A,0xED), "acento": RGBColor(0xFE,0xF9,0xC3), "texto": RGBColor(0x1C,0x19,0x17), "texto_sub": RGBColor(0x57,0x53,0x4E)},
    "corporativo":  {"fundo": RGBColor(0x1E,0x3A,0x5F), "primaria": RGBColor(0xC0,0x39,0x2B), "secundaria": RGBColor(0xEC,0xF0,0xF1), "acento": RGBColor(0xF3,0x9C,0x12), "texto": RGBColor(0xFF,0xFF,0xFF), "texto_sub": RGBColor(0xBD,0xC3,0xC7)},
    "saude":        {"fundo": RGBColor(0xF0,0xFF,0xF4), "primaria": RGBColor(0x05,0x7A,0x55), "secundaria": RGBColor(0x06,0x95,0xDD), "acento": RGBColor(0xD1,0xFA,0xE5), "texto": RGBColor(0x06,0x4E,0x3B), "texto_sub": RGBColor(0x06,0x5F,0x46)},
    "tech":         {"fundo": RGBColor(0x0D,0x11,0x17), "primaria": RGBColor(0x00,0xFF,0xC8), "secundaria": RGBColor(0x7C,0x3A,0xED), "acento": RGBColor(0x1A,0x1F,0x2E), "texto": RGBColor(0xFF,0xFF,0xFF), "texto_sub": RGBColor(0x94,0xA3,0xB8)},
}

def _set_bg(slide, cor):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = cor

def _rect(slide, l, t, w, h, cor):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = cor; s.line.fill.background(); return s

def _txt(slide, texto, l, t, w, h, size=18, cor=RGBColor(255,255,255), bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = texto
    r.font.size = Pt(size); r.font.color.rgb = cor
    r.font.bold = bold; r.font.italic = italic; r.font.name = "Segoe UI"
    return tb


def gerar_pptx_com_imagens(
    titulo: str,
    subtitulo: str,
    slides_conteudo: list,
    categoria: str = "apresentacao",
    gerar_imagens: bool = True,
    tema_override: str = None
) -> str:
    """
    Gera PPTX profissional com imagens contextuais geradas por IA.

    Args:
        titulo: Título da apresentação
        subtitulo: Subtítulo
        slides_conteudo: Lista de dicts [{"titulo": "...", "conteudo": "...", "tipo": "texto|lista"}]
        categoria: Categoria do template (define paleta de cores e estilo visual)
        gerar_imagens: Se True, gera imagens com Gemini + DALL-E
        tema_override: Se fornecido, usa esse tema como contexto para as imagens

    Returns:
        Caminho do arquivo PPTX gerado
    """
    paleta = PALETAS.get(categoria, PALETAS["apresentacao"])
    tema_imagem = tema_override or titulo

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    # ── SLIDE DE CAPA com imagem contextual ──
    s_capa = prs.slides.add_slide(layout)
    _set_bg(s_capa, paleta["fundo"])

    # Gerar imagem de capa
    if gerar_imagens:
        try:
            print(f"🎨 Gerando imagem de capa para: '{tema_imagem}'")
            img_capa = gerar_imagem_contextual(
                tema=tema_imagem,
                categoria=categoria,
                tipo_imagem="capa de apresentação profissional, wide banner",
                tamanho="horizontal"
            )
            # Inserir imagem como fundo semitransparente (lado direito)
            s_capa.shapes.add_picture(img_capa, Inches(4.5), Inches(0), Inches(5.5), Inches(7.5))
        except Exception as e:
            print(f"⚠️ Imagem de capa não gerada: {e}")

    # Overlay escuro sobre a imagem para legibilidade
    overlay = _rect(s_capa, 4.5, 0, 5.5, 7.5, paleta["fundo"])
    overlay.fill.fore_color.rgb = paleta["fundo"]
    # Transparência via XML
    from pptx.oxml.ns import qn
    solidFill = overlay.fill._xPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is None:
            from lxml import etree
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            r, g, b = paleta["fundo"]
            srgbClr.set('val', f'{r:02X}{g:02X}{b:02X}')
        alpha = srgbClr.find(qn('a:alpha'))
        if alpha is None:
            from lxml import etree
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '80000')  # ~80% opacidade

    # Elementos da capa
    _rect(s_capa, 0, 0, 0.12, 7.5, paleta["primaria"])
    _rect(s_capa, 0, 6.5, 10, 1.0, paleta["primaria"])
    _rect(s_capa, 0.3, 3.6, 6, 0.05, paleta["secundaria"])
    _txt(s_capa, titulo.upper(), 0.3, 1.5, 9, 1.5, size=34, cor=paleta["texto"], bold=True)
    _txt(s_capa, subtitulo, 0.3, 3.8, 9, 0.7, size=17, cor=paleta["texto_sub"], italic=True)
    from datetime import datetime
    _txt(s_capa, f"Sistema IA  |  {datetime.now().strftime('%d/%m/%Y')}", 0.3, 6.55, 9, 0.4,
         size=12, cor=paleta["texto"], align=PP_ALIGN.CENTER)

    # ── SLIDES DE CONTEÚDO ──
    for idx, slide_data in enumerate(slides_conteudo):
        tipo = slide_data.get("tipo", "texto")
        titulo_s = slide_data.get("titulo", "")
        conteudo_s = slide_data.get("conteudo", "")
        gerar_img_slide = slide_data.get("imagem", False)  # Opcional por slide

        s = prs.slides.add_slide(layout)
        _set_bg(s, paleta["fundo"])
        _rect(s, 0, 0, 10, 1.3, paleta["primaria"])
        _rect(s, 0, 0, 0.12, 7.5, paleta["secundaria"])
        _txt(s, titulo_s.upper() if titulo_s else "", 0.3, 0.2, 9.5, 0.9,
             size=24, cor=paleta["texto"], bold=True)
        _rect(s, 0.3, 1.4, 9.4, 0.04, paleta["secundaria"])

        # Gerar imagem contextual para o slide (se solicitado)
        if gerar_imagens and gerar_img_slide and titulo_s:
            try:
                print(f"🎨 Gerando imagem para slide: '{titulo_s}'")
                img_slide = gerar_imagem_contextual(
                    tema=f"{titulo_s} - {tema_imagem}",
                    categoria=categoria,
                    tipo_imagem="ilustração de slide de apresentação",
                    tamanho="quadrado"
                )
                # Imagem no lado direito do slide
                s.shapes.add_picture(img_slide, Inches(6.5), Inches(1.6), Inches(3.2), Inches(5.5))
                largura_texto = 5.8  # Texto ocupa lado esquerdo
            except Exception as e:
                print(f"⚠️ Imagem do slide não gerada: {e}")
                largura_texto = 9.0
        else:
            largura_texto = 9.0

        if tipo == "lista" and isinstance(conteudo_s, list):
            y = 1.6
            for item in conteudo_s[:6]:
                _rect(s, 0.5, y + 0.1, 0.22, 0.22, paleta["acento"])
                _txt(s, str(item), 0.9, y, largura_texto - 0.9, 0.55, size=15, cor=paleta["texto"])
                y += 0.72
        else:
            _txt(s, str(conteudo_s), 0.5, 1.6, largura_texto, 5.5, size=16, cor=paleta["texto"])

    # ── SLIDE DE ENCERRAMENTO ──
    s_fim = prs.slides.add_slide(layout)
    _set_bg(s_fim, paleta["primaria"])
    _rect(s_fim, 0, 0, 10, 0.08, paleta["secundaria"])
    _rect(s_fim, 0, 7.42, 10, 0.08, paleta["secundaria"])
    _txt(s_fim, "OBRIGADO!", 0, 2.8, 10, 1.2, size=44, cor=paleta["texto"], bold=True, align=PP_ALIGN.CENTER)
    _txt(s_fim, "Dúvidas? Entre em contato.", 0, 4.2, 10, 0.6, size=18,
         cor=paleta["texto_sub"], align=PP_ALIGN.CENTER, italic=True)

    # Salvar
    uid = str(uuid.uuid4())[:8]
    caminho = f"outputs/{titulo.replace(' ', '_')}_{uid}.pptx"
    prs.save(caminho)
    print(f"\n✅ PPTX com imagens gerado: {caminho}")
    return caminho


# ── Exemplo de uso ──
if __name__ == "__main__":
    slides = [
        {
            "titulo": "Visão Geral do Mercado",
            "conteudo": "O mercado de IA cresceu 145% em 2025, com projeção de US$ 1.8 trilhão até 2030.",
            "tipo": "texto",
            "imagem": True   # ← gera imagem contextual para este slide
        },
        {
            "titulo": "Principais Oportunidades",
            "conteudo": [
                "Automação de processos com agentes de IA",
                "Análise preditiva e business intelligence",
                "Geração de conteúdo personalizado",
                "Atendimento ao cliente 24/7",
            ],
            "tipo": "lista",
            "imagem": False
        },
        {
            "titulo": "Resultados Esperados",
            "conteudo": "Redução de 40% nos custos operacionais e aumento de 60% na produtividade.",
            "tipo": "texto",
            "imagem": True
        },
    ]

    gerar_pptx_com_imagens(
        titulo="Estratégia de IA para 2026",
        subtitulo="Oportunidades e Resultados",
        slides_conteudo=slides,
        categoria="tech",
        gerar_imagens=True
    )
