"""
CRIADOR DE TEMPLATES PPTX PROFISSIONAIS
Categorias: Financeiro, Marketing, Apresentação, Escola, Ideias, Corporativo, Saúde, Tech
Uso: python criar_templates_pptx.py
Resultado: arquivos .pptx na pasta templates/pptx/
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

os.makedirs("templates/pptx", exist_ok=True)

# ── Paletas de cores por categoria ──
PALETAS = {
    "financeiro": {
        "fundo": RGBColor(0x0A, 0x1F, 0x44),       # Azul marinho profundo
        "primaria": RGBColor(0x00, 0x8B, 0x8B),     # Ciano escuro
        "secundaria": RGBColor(0x00, 0xD4, 0xAA),   # Verde água
        "acento": RGBColor(0xF0, 0xB4, 0x29),       # Dourado
        "texto": RGBColor(0xFF, 0xFF, 0xFF),
        "texto_sub": RGBColor(0xB0, 0xC4, 0xDE),
    },
    "marketing": {
        "fundo": RGBColor(0x1A, 0x00, 0x2E),        # Roxo escuro
        "primaria": RGBColor(0xFF, 0x00, 0x80),     # Rosa neon
        "secundaria": RGBColor(0xFF, 0x6B, 0x35),   # Laranja vibrante
        "acento": RGBColor(0xFF, 0xD7, 0x00),       # Amarelo
        "texto": RGBColor(0xFF, 0xFF, 0xFF),
        "texto_sub": RGBColor(0xE0, 0xC0, 0xFF),
    },
    "apresentacao": {
        "fundo": RGBColor(0xFF, 0xFF, 0xFF),        # Branco
        "primaria": RGBColor(0x2C, 0x3E, 0x50),     # Cinza azulado escuro
        "secundaria": RGBColor(0x3B, 0x82, 0xF6),   # Azul moderno
        "acento": RGBColor(0xE8, 0xF4, 0xFD),       # Azul muito claro
        "texto": RGBColor(0x1E, 0x29, 0x3B),
        "texto_sub": RGBColor(0x64, 0x74, 0x8B),
    },
    "escola": {
        "fundo": RGBColor(0xF0, 0xF9, 0xFF),        # Azul clarinho
        "primaria": RGBColor(0x05, 0x96, 0x68),     # Verde educação
        "secundaria": RGBColor(0xF5, 0x9E, 0x0B),   # Laranja amigável
        "acento": RGBColor(0xFE, 0xF3, 0xC7),       # Amarelo suave
        "texto": RGBColor(0x1E, 0x29, 0x3B),
        "texto_sub": RGBColor(0x37, 0x41, 0x51),
    },
    "ideias": {
        "fundo": RGBColor(0xFF, 0xFD, 0xF0),        # Creme quente
        "primaria": RGBColor(0xD9, 0x77, 0x06),     # Âmbar
        "secundaria": RGBColor(0x7C, 0x3A, 0xED),   # Roxo criativo
        "acento": RGBColor(0xFE, 0xF9, 0xC3),       # Amarelo post-it
        "texto": RGBColor(0x1C, 0x19, 0x17),
        "texto_sub": RGBColor(0x57, 0x53, 0x4E),
    },
    "corporativo": {
        "fundo": RGBColor(0x1E, 0x3A, 0x5F),        # Azul corporativo
        "primaria": RGBColor(0xC0, 0x39, 0x2B),     # Vermelho corporativo
        "secundaria": RGBColor(0xEC, 0xF0, 0xF1),   # Cinza claro
        "acento": RGBColor(0xF3, 0x9C, 0x12),       # Laranja
        "texto": RGBColor(0xFF, 0xFF, 0xFF),
        "texto_sub": RGBColor(0xBD, 0xC3, 0xC7),
    },
    "saude": {
        "fundo": RGBColor(0xF0, 0xFF, 0xF4),        # Verde muito claro
        "primaria": RGBColor(0x05, 0x7A, 0x55),     # Verde saúde
        "secundaria": RGBColor(0x06, 0x95, 0xDD),   # Azul médico
        "acento": RGBColor(0xD1, 0xFA, 0xE5),       # Verde suave
        "texto": RGBColor(0x06, 0x4E, 0x3B),
        "texto_sub": RGBColor(0x06, 0x5F, 0x46),
    },
    "tech": {
        "fundo": RGBColor(0x0D, 0x11, 0x17),        # Preto tech
        "primaria": RGBColor(0x00, 0xFF, 0xC8),     # Ciano neon
        "secundaria": RGBColor(0x7C, 0x3A, 0xED),   # Roxo tech
        "acento": RGBColor(0x1A, 0x1F, 0x2E),       # Cinza escuro
        "texto": RGBColor(0xFF, 0xFF, 0xFF),
        "texto_sub": RGBColor(0x94, 0xA3, 0xB8),
    },
}

def _set_background(slide, cor: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = cor

def _add_rect(slide, left, top, width, height, cor: RGBColor):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor
    shape.line.fill.background()
    return shape

def _add_text(slide, texto, left, top, width, height,
              size=20, cor=RGBColor(255,255,255), bold=False,
              align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(size)
    run.font.color.rgb = cor
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Segoe UI"
    return txBox

def criar_template_pptx(nome: str, paleta: dict, icone_categoria: str = "●"):
    """Cria um template PPTX completo com 4 layouts de slide."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    layout = prs.slide_layouts[6]  # Blank

    # ── SLIDE 1: Capa ──
    s1 = prs.slides.add_slide(layout)
    _set_background(s1, paleta["fundo"])

    # Barra lateral esquerda
    _add_rect(s1, 0, 0, 0.12, 7.5, paleta["primaria"])
    # Barra inferior
    _add_rect(s1, 0, 6.5, 10, 1.0, paleta["primaria"])
    # Linha decorativa
    _add_rect(s1, 0.3, 3.6, 7, 0.05, paleta["secundaria"])

    _add_text(s1, icone_categoria, 0.3, 0.5, 1, 0.8, size=36, cor=paleta["secundaria"], bold=True)
    _add_text(s1, "TÍTULO DA APRESENTAÇÃO", 0.3, 1.5, 9, 1.2, size=36, cor=paleta["texto"], bold=True)
    _add_text(s1, "Subtítulo ou descrição aqui", 0.3, 3.8, 9, 0.6, size=18, cor=paleta["texto_sub"], italic=True)
    _add_text(s1, "Autor  |  Data  |  Versão", 0.3, 6.55, 9, 0.4, size=12, cor=paleta["texto"], align=PP_ALIGN.CENTER)

    # ── SLIDE 2: Conteúdo com texto ──
    s2 = prs.slides.add_slide(layout)
    _set_background(s2, paleta["fundo"])
    _add_rect(s2, 0, 0, 10, 1.3, paleta["primaria"])
    _add_rect(s2, 0, 0, 0.12, 7.5, paleta["secundaria"])
    _add_rect(s2, 0.3, 1.4, 9.4, 0.04, paleta["secundaria"])

    _add_text(s2, "TÍTULO DO SLIDE", 0.3, 0.2, 9.5, 0.9, size=26, cor=paleta["texto"], bold=True)
    _add_text(s2, "Insira o conteúdo do slide aqui. Este é o layout padrão para texto corrido, explicações e descrições detalhadas.", 0.5, 1.6, 9, 5.0, size=17, cor=paleta["texto"])

    # ── SLIDE 3: Lista de tópicos ──
    s3 = prs.slides.add_slide(layout)
    _set_background(s3, paleta["fundo"])
    _add_rect(s3, 0, 0, 10, 1.3, paleta["primaria"])
    _add_rect(s3, 0, 0, 0.12, 7.5, paleta["secundaria"])

    _add_text(s3, "TÓPICOS PRINCIPAIS", 0.3, 0.2, 9.5, 0.9, size=26, cor=paleta["texto"], bold=True)

    topicos = ["Primeiro tópico importante", "Segundo tópico relevante",
               "Terceiro ponto de destaque", "Quarto item da lista",
               "Quinto elemento da apresentação"]
    y = 1.5
    for t in topicos:
        _add_rect(s3, 0.5, y + 0.1, 0.22, 0.22, paleta["acento"])
        _add_text(s3, t, 0.9, y, 8.8, 0.55, size=16, cor=paleta["texto"])
        y += 0.72

    # ── SLIDE 4: Dois colunas ──
    s4 = prs.slides.add_slide(layout)
    _set_background(s4, paleta["fundo"])
    _add_rect(s4, 0, 0, 10, 1.3, paleta["primaria"])
    _add_rect(s4, 0, 0, 0.12, 7.5, paleta["secundaria"])
    _add_rect(s4, 5.1, 1.4, 0.05, 5.8, paleta["secundaria"])  # Divisor central

    _add_text(s4, "COMPARATIVO / DOIS TEMAS", 0.3, 0.2, 9.5, 0.9, size=26, cor=paleta["texto"], bold=True)
    _add_text(s4, "Coluna Esquerda", 0.5, 1.5, 4.4, 0.6, size=16, cor=paleta["acento"], bold=True)
    _add_text(s4, "Insira o conteúdo da coluna esquerda aqui.", 0.5, 2.2, 4.4, 4.5, size=14, cor=paleta["texto"])
    _add_text(s4, "Coluna Direita", 5.3, 1.5, 4.4, 0.6, size=16, cor=paleta["acento"], bold=True)
    _add_text(s4, "Insira o conteúdo da coluna direita aqui.", 5.3, 2.2, 4.4, 4.5, size=14, cor=paleta["texto"])

    # ── SLIDE 5: Encerramento ──
    s5 = prs.slides.add_slide(layout)
    _set_background(s5, paleta["primaria"])
    _add_rect(s5, 0, 0, 10, 0.08, paleta["secundaria"])
    _add_rect(s5, 0, 7.42, 10, 0.08, paleta["secundaria"])

    _add_text(s5, icone_categoria, 4.2, 1.5, 2, 1, size=48, cor=paleta["acento"], bold=True, align=PP_ALIGN.CENTER)
    _add_text(s5, "OBRIGADO!", 0, 2.8, 10, 1.2, size=44, cor=paleta["texto"], bold=True, align=PP_ALIGN.CENTER)
    _add_text(s5, "Dúvidas? Entre em contato.", 0, 4.2, 10, 0.6, size=18, cor=paleta["texto_sub"], align=PP_ALIGN.CENTER, italic=True)
    _add_text(s5, "contato@empresa.com  |  www.empresa.com", 0, 5.5, 10, 0.5, size=13, cor=paleta["texto_sub"], align=PP_ALIGN.CENTER)

    caminho = f"templates/pptx/template_{nome}.pptx"
    prs.save(caminho)
    print(f"✅ Template PPTX criado: {caminho}")
    return caminho


# ── Criar todos os templates ──
if __name__ == "__main__":
    configs = [
        ("financeiro",    PALETAS["financeiro"],    "💰"),
        ("marketing",     PALETAS["marketing"],     "📣"),
        ("apresentacao",  PALETAS["apresentacao"],  "📊"),
        ("escola",        PALETAS["escola"],        "📚"),
        ("ideias",        PALETAS["ideias"],        "💡"),
        ("corporativo",   PALETAS["corporativo"],   "🏢"),
        ("saude",         PALETAS["saude"],         "🏥"),
        ("tech",          PALETAS["tech"],          "⚙"),
    ]

    for nome, paleta, icone in configs:
        criar_template_pptx(nome, paleta, icone)

    print("\n✅ Todos os templates PPTX foram criados em templates/pptx/")
