import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from fpdf import FPDF

# --- CONFIGURAÇÕES E TEMPLATES ---

# Cores da Marca (Exemplo)
COR_PRIMARIA = (31, 73, 125) # Azul Escuro
COR_SECUNDARIA = (255, 255, 255) # Branco
COR_TEXTO = (0, 0, 0) # Preto

# --- FUNÇÕES DE GERAÇÃO DE ARQUIVOS (A LÓGICA DE FORMATAÇÃO) ---

def criar_pptx_profissional(dados_conteudo, nome_arquivo="apresentacao_profissional.pptx"):
    """
    Gera um arquivo PPTX profissional com base em dados estruturados.
    A IA fornece os dados (dados_conteudo), mas o código controla a formatação.
    """
    prs = Presentation()

    # Slide de Título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = dados_conteudo.get('titulo', 'Título da Apresentação')
    subtitle.text = dados_conteudo.get('subtitulo', 'Subtítulo ou Descrição Breve')

    # Slides de Conteúdo
    for item in dados_conteudo.get('slides', []):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = item.get('titulo_slide', 'Título do Slide')
        
        tf = body_shape.text_frame
        tf.text = item.get('texto_principal', '')

        for ponto in item.get('pontos_chave', []):
            p = tf.add_paragraph()
            p.text = ponto
            p.level = 1

    prs.save(nome_arquivo)
    print(f"PPTX '{nome_arquivo}' gerado com sucesso.")

def criar_pdf_profissional(dados_conteudo, nome_arquivo="relatorio_profissional.pdf"):
    """
    Gera um arquivo PDF profissional com base em dados estruturados.
    """
    pdf = FPDF()
    pdf.add_page()
    
    # Cabeçalho
    pdf.set_fill_color(*COR_PRIMARIA)
    pdf.rect(0, 0, 210, 40, 'F')
    
    pdf.set_font('Arial', 'B', 24)
    pdf.set_text_color(*COR_SECUNDARIA)
    pdf.cell(0, 20, dados_conteudo.get('titulo', 'Relatório Profissional'), ln=True, align='C')
    
    pdf.set_font('Arial', 'I', 12)
    pdf.cell(0, 10, dados_conteudo.get('subtitulo', 'Subtítulo do Relatório'), ln=True, align='C')
    
    pdf.ln(20)
    pdf.set_text_color(*COR_TEXTO)
    
    # Conteúdo
    for secao in dados_conteudo.get('secoes', []):
        pdf.set_font('Arial', 'B', 16)
        pdf.cell(0, 10, secao.get('titulo_secao', 'Título da Seção'), ln=True)
        pdf.ln(5)
        
        pdf.set_font('Arial', '', 12)
        pdf.multi_cell(0, 10, secao.get('texto_secao', 'Conteúdo da seção...'))
        pdf.ln(10)

    pdf.output(nome_arquivo)
    print(f"PDF '{nome_arquivo}' gerado com sucesso.")

# --- EXEMPLO DE USO (O QUE O SEU AGENTE DE IA FARIA) ---

if __name__ == "__main__":
    # 1. O Agente de IA gera apenas este dicionário de dados (JSON/Dicionário)
    # Ele NÃO gera o código de formatação acima.
    conteudo_exemplo = {
        'titulo': 'Inovação com IA em 2026',
        'subtitulo': 'Tendências e Aplicações Práticas para Empresas',
        'slides': [
            {
                'titulo_slide': 'O Estado Atual da IA',
                'texto_principal': 'A IA evoluiu de modelos de linguagem para agentes autônomos capazes de realizar tarefas complexas.',
                'pontos_chave': [
                    'Agentes que operam ferramentas e navegadores.',
                    'Multimodalidade nativa (texto, imagem, áudio, vídeo).',
                    'Integração profunda em fluxos de trabalho empresariais.'
                ]
            },
            {
                'titulo_slide': 'Benefícios da Automação',
                'texto_principal': 'Empresas que adotam agentes de IA veem ganhos significativos em eficiência e qualidade.',
                'pontos_chave': [
                    'Redução de erros humanos em tarefas repetitivas.',
                    'Geração de documentos e relatórios em segundos.',
                    'Análise de dados em tempo real para tomada de decisão.'
                ]
            }
        ],
        'secoes': [
            {
                'titulo_secao': 'Introdução à Nova Era da IA',
                'texto_secao': 'Estamos vivendo uma mudança de paradigma onde a IA não é apenas uma ferramenta de consulta, mas um colaborador ativo que entende o contexto e executa ações complexas em nome do usuário.'
            },
            {
                'titulo_secao': 'Conclusão e Próximos Passos',
                'texto_secao': 'Para implementar essas soluções com sucesso, é fundamental focar na qualidade da integração e na clareza dos fluxos de trabalho entre humanos e máquinas.'
            }
        ]
    }

    # 2. O seu sistema chama as funções de geração passando os dados
    criar_pptx_profissional(conteudo_exemplo)
    criar_pdf_profissional(conteudo_exemplo)
