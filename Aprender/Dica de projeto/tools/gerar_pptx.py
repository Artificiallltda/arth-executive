"""
GERADOR DE PPTX DE ALTA QUALIDADE
Usa: python-pptx com design profissional
Resultado: Apresentação com layout, cores e tipografia profissionais
"""

import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# ── Temas de cores ──
TEMAS_PPTX = {
    "profissional": {
        "fundo": RGBColor(0x0F, 0x17, 0x2A),       # Azul muito escuro
        "primaria": RGBColor(0x3B, 0x82, 0xF6),     # Azul
        "secundaria": RGBColor(0x06, 0xB6, 0xD4),   # Ciano
        "texto": RGBColor(0xF8, 0xFA, 0xFC),        # Branco
        "texto_sub": RGBColor(0x94, 0xA3, 0xB8),    # Cinza claro
        "acento": RGBColor(0x1E, 0x40, 0xAF),       # Azul médio
    },
    "moderno": {
        "fundo": RGBColor(0xFF, 0xFF, 0xFF),
        "primaria": RGBColor(0x6D, 0x28, 0xD9),
        "secundaria": RGBColor(0xEC, 0x48, 0x99),
        "texto": RGBColor(0x1E, 0x29, 0x3B),
        "texto_sub": RGBColor(0x64, 0x74, 0x8B),
        "acento": RGBColor(0xED, 0xE9, 0xFE),
    },
    "corporativo": {
        "fundo": RGBColor(0xFF, 0xFF, 0xFF),
        "primaria": RGBColor(0x1E, 0x3A, 0x8A),
        "secundaria": RGBColor(0x16, 0xA3, 0x4A),
        "texto": RGBColor(0x1E, 0x29, 0x3B),
        "texto_sub": RGBColor(0x64, 0x74, 0x8B),
        "acento": RGBColor(0xDB, 0xEA, 0xFE),
    },
}

def _set_background(slide, cor: RGBColor):
    """Define cor de fundo do slide."""
    from pptx.oxml.ns import qn
    from lxml import etree
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = cor

def _adicionar_retangulo(slide, left, top, width, height, cor: RGBColor, transparencia=0):
    """Adiciona retângulo colorido ao slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor
    shape.line.fill.background()
    return shape

def _adicionar_texto(slide, texto, left, top, width, height, 
                     tamanho=24, cor=RGBColor(0xFF,0xFF,0xFF), 
                     negrito=False, alinhamento=PP_ALIGN.LEFT, italico=False):
    """Adiciona caixa de texto formatada."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alinhamento
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.color.rgb = cor
    run.font.bold = negrito
    run.font.italic = italico
    run.font.name = "Segoe UI"
    return txBox

def _slide_capa(prs, titulo, subtitulo, tema):
    """Cria slide de capa profissional."""
    slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(slide_layout)
    
    _set_background(slide, tema["fundo"])
    
    # Barra lateral decorativa
    _adicionar_retangulo(slide, 0, 0, 0.08, 7.5, tema["primaria"])
    
    # Linha decorativa horizontal
    _adicionar_retangulo(slide, 0.5, 3.8, 8, 0.04, tema["secundaria"])
    
    # Título principal
    _adicionar_texto(slide, titulo, 0.7, 1.5, 8.5, 1.8,
                     tamanho=44, cor=tema["texto"], negrito=True)
    
    # Subtítulo
    if subtitulo:
        _adicionar_texto(slide, subtitulo, 0.7, 4.0, 8.5, 0.8,
                         tamanho=20, cor=tema["texto_sub"], italico=True)
    
    # Linha decorativa inferior
    _adicionar_retangulo(slide, 0, 6.8, 10, 0.7, tema["acento"])
    _adicionar_texto(slide, "Apresentação Profissional", 0.3, 6.85, 9, 0.5,
                     tamanho=12, cor=tema["texto_sub"], alinhamento=PP_ALIGN.CENTER)
    
    return slide

def _slide_conteudo_texto(prs, titulo_slide, conteudo, tema):
    """Cria slide de conteúdo com texto."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    _set_background(slide, tema["fundo"])
    
    # Cabeçalho colorido
    _adicionar_retangulo(slide, 0, 0, 10, 1.4, tema["acento"])
    _adicionar_retangulo(slide, 0, 0, 0.08, 1.4, tema["primaria"])
    
    # Título do slide
    _adicionar_texto(slide, titulo_slide, 0.3, 0.2, 9.5, 1.0,
                     tamanho=28, cor=tema["primaria"], negrito=True)
    
    # Linha separadora
    _adicionar_retangulo(slide, 0.3, 1.5, 9.4, 0.03, tema["secundaria"])
    
    # Conteúdo
    _adicionar_texto(slide, conteudo, 0.5, 1.7, 9, 5.0,
                     tamanho=18, cor=tema["texto"])
    
    return slide

def _slide_lista(prs, titulo_slide, itens: list, tema):
    """Cria slide com lista de tópicos."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    _set_background(slide, tema["fundo"])
    
    # Cabeçalho
    _adicionar_retangulo(slide, 0, 0, 10, 1.4, tema["acento"])
    _adicionar_retangulo(slide, 0, 0, 0.08, 1.4, tema["primaria"])
    _adicionar_texto(slide, titulo_slide, 0.3, 0.2, 9.5, 1.0,
                     tamanho=28, cor=tema["primaria"], negrito=True)
    
    # Itens da lista
    y_pos = 1.7
    for i, item in enumerate(itens[:6]):  # Máximo 6 itens por slide
        # Marcador colorido
        _adicionar_retangulo(slide, 0.5, y_pos + 0.12, 0.25, 0.25, tema["secundaria"])
        # Texto do item
        _adicionar_texto(slide, item, 0.9, y_pos, 8.8, 0.6,
                         tamanho=17, cor=tema["texto"])
        y_pos += 0.75
    
    return slide

def gerar_pptx(titulo_apresentacao: str, subtitulo: str, slides: list, tema_nome: str = "profissional") -> str:
    """
    Gera PPTX profissional.
    
    slides: lista de dicts com:
      - titulo: str
      - conteudo: str (texto) ou list (para lista de tópicos)
      - tipo: "texto" | "lista"
    """
    tema = TEMAS_PPTX.get(tema_nome, TEMAS_PPTX["profissional"])
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide de capa
    _slide_capa(prs, titulo_apresentacao, subtitulo, tema)
    
    # Slides de conteúdo
    for slide_data in slides:
        tipo = slide_data.get("tipo", "texto")
        titulo_s = slide_data.get("titulo", "")
        conteudo_s = slide_data.get("conteudo", "")
        
        if tipo == "lista" and isinstance(conteudo_s, list):
            _slide_lista(prs, titulo_s, conteudo_s, tema)
        else:
            _slide_conteudo_texto(prs, titulo_s, str(conteudo_s), tema)
    
    # Salvar
    os.makedirs("outputs", exist_ok=True)
    uid = str(uuid.uuid4())[:8]
    caminho = f"outputs/{titulo_apresentacao.replace(' ', '_')}_{uid}.pptx"
    prs.save(caminho)
    print(f"✅ PPTX gerado: {caminho}")
    return caminho


# Uso direto
if __name__ == "__main__":
    slides_exemplo = [
        {
            "titulo": "O que é Inteligência Artificial?",
            "conteudo": "A Inteligência Artificial é a simulação de processos de inteligência humana por sistemas computacionais, incluindo aprendizado, raciocínio e autocorreção.",
            "tipo": "texto"
        },
        {
            "titulo": "Principais Aplicações",
            "conteudo": [
                "Assistentes virtuais e chatbots",
                "Análise preditiva e business intelligence",
                "Geração de conteúdo e imagens",
                "Automação de processos (RPA + IA)",
                "Diagnóstico médico e saúde",
                "Sistemas de recomendação"
            ],
            "tipo": "lista"
        },
        {
            "titulo": "Tendências para 2026",
            "conteudo": "Os squads de agentes de IA representam a próxima fronteira da automação inteligente, com crescimento de 1.445% na demanda por sistemas multiagentes.",
            "tipo": "texto"
        },
    ]
    
    gerar_pptx(
        titulo_apresentacao="Inteligência Artificial em 2026",
        subtitulo="Tendências, Aplicações e Oportunidades",
        slides=slides_exemplo,
        tema_nome="profissional"
    )
