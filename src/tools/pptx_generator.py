# 🛡️ SKILL BLINDADA (14/03/2026) - PADRÃO MANUS AI (GIT PERSISTENCE)
import os
import json
import uuid
import logging
import re
import asyncio
from langchain_core.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from src.config import settings
from src.tools.image_generator import generate_image
from typing import Any, Optional
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

# ─── CONFIGURAÇÃO DE CAMINHOS ABSOLUTOS (MANUS STYLE) ───────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__)) # src/tools
PROJECT_ROOT = os.path.dirname(os.path.dirname(BASE_DIR)) # arth-executive
TEMPLATES_DIR = os.getenv("TEMPLATES_DIR", os.path.join(PROJECT_ROOT, "data", "templates", "pptx"))

def _get_template_path(categoria: str = None):
    """Busca um template de forma ABSOLUTA com Fallback (Manus AI)."""
    logger.info(f"🔍 [PPTXGen] Buscando template para categoria: '{categoria}' em {TEMPLATES_DIR}")
    
    # 1. Caminho do template desejado
    cat_name = str(categoria or "apresentacao").lower().replace("template_", "").replace(".pptx", "")
    template_path = os.path.join(TEMPLATES_DIR, f"template_{cat_name}.pptx")

    # 2. Verificação e Fallback Automático
    if not os.path.exists(template_path):
        logger.warning(f"⚠️ [PPTXGen] Template '{template_path}' não encontrado. Usando fallback.")
        template_path = os.path.join(TEMPLATES_DIR, "template_apresentacao.pptx")

    # 3. Fallback Final (Criação do zero se nada existir)
    if not os.path.exists(template_path):
        logger.error(f"❌ [PPTXGen] Nenhum template encontrado em {template_path}. Gerando sem template.")
        return None
    
    logger.info(f"✅ [PPTXGen] Usando template absoluto: {template_path}")
    return template_path

# ─── Manus Executive Design System (Fallback Code) ─────────────────────────
_BG, _CARD, _ACCENT = RGBColor(10, 12, 16), RGBColor(28, 33, 40), RGBColor(88, 166, 255)
_WHITE, _TEXT = RGBColor(255, 255, 255), RGBColor(230, 237, 243)
_W, _H = Inches(13.333), Inches(7.5)

def _bg(slide, color=_BG):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color

def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def _text(slide, l, t, w, h, txt, size, bold=False, color=_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(txt); r.font.name = "Calibri"
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

def _build_cover(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide); _rect(slide, Inches(0), Inches(2.2), _W, Inches(3.2), _CARD)
    font_size = 54 if len(title) < 30 else 42
    _text(slide, Inches(1.2), Inches(2.5), Inches(11.0), Inches(1.8), title.upper(), font_size, bold=True, color=_WHITE)
    if subtitle: _text(slide, Inches(1.2), Inches(4.3), Inches(10.5), Inches(0.8), subtitle, 22, color=_ACCENT)

def _build_content(prs, title, bullets, img_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide); _rect(slide, Inches(0), Inches(0), _W, Inches(1.1), _CARD)
    _text(slide, Inches(0.75), Inches(0.30), Inches(12), Inches(0.6), title.upper(), 28, bold=True, color=_WHITE)
    
    # Injeção de imagem simplificada
    has_image = False
    if img_path:
        clean_name = re.sub(r'<[^>]+>', '', str(img_path)).replace("SEND_FILE:", "").strip()
        fp = os.path.join(settings.DATA_OUTPUTS_PATH, clean_name)
        if os.path.exists(fp):
            try:
                slide.shapes.add_picture(fp, Inches(0.5), Inches(1.5), width=Inches(6), height=Inches(5))
                has_image = True
            except: pass

    l, w = (Inches(7.0), Inches(5.8)) if has_image else (Inches(1.0), Inches(11.3))
    tb = slide.shapes.add_textbox(l, Inches(1.6), w, Inches(5.3))
    tf = tb.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = f"◈ {bullet}"; r.font.size = Pt(19); r.font.color.rgb = _TEXT

class PpptxSchema(BaseModel):
    slides_content_json: Any = Field(..., description="Conteúdo dos slides em JSON.")
    template_name: Optional[str] = Field(None, description="Categoria do template (ex: financeiro, tech).")

@tool(args_schema=PpptxSchema)
async def generate_pptx(slides_content_json: Any, template_name: str = None) -> str:
    """Gera PPTX Premium com suporte a TEMPLATES PERSISTENTES."""
    try:
        filename = f"Exec-Deck-{uuid.uuid4().hex[:6]}.pptx"
        filepath = os.path.join(settings.DATA_OUTPUTS_PATH, filename)

        # Normalização de JSON
        prs_title, prs_subtitle, slides_data = "EXECUTIVE DECK", "", []
        if isinstance(slides_content_json, str):
            clean = re.sub(r'```(?:json)?\n?(.*?)\n?```', r'\1', slides_content_json, flags=re.DOTALL).strip()
            slides_content_json = json.loads(clean)
        
        if isinstance(slides_content_json, dict):
            prs_title = slides_content_json.get("title", "Apresentação")
            prs_subtitle = slides_content_json.get("subtitle", "")
            slides_data = slides_content_json.get("slides", [])

        # Carregamento do Template Blindado
        t_path = _get_template_path(template_name)
        prs = Presentation(t_path) if t_path else Presentation()
        if not t_path: prs.slide_width, prs.slide_height = _W, _H

        _build_cover(prs, prs_title, prs_subtitle)
        for i, s in enumerate(slides_data):
            _build_content(prs, s.get("title", f"SLIDE {i+1}"), s.get("bullets", []), s.get("image"))

        await asyncio.to_thread(prs.save, filepath)
        return f"Apresentação gerada com sucesso: <SEND_FILE:{filename}>"
    except Exception as e:
        logger.error(f"[PPTX] Erro: {e}")
        return f"Falha no PPTX: {str(e)}"
